"""
Presentation Service — NotebookLM Technical Blueprint & Visual Slide Engine
Generates professional PowerPoint presentations (.pptx) and interactive HTML slide decks (Reveal.js)
from structured JSON slide specifications, supporting embedded AI images, process pipelines, 
blueprint diagrams, and split layouts.
"""

import os
import json
import logging
from typing import Dict, Any, List, Optional

logger = logging.getLogger(__name__)


def create_powerpoint_deck(spec: Dict[str, Any], output_path: str) -> Dict[str, Any]:
    """
    Build a high-end, executive-grade PowerPoint presentation (.pptx) in NotebookLM Technical Blueprint Style
    from a JSON slide specification.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title = spec.get("title", "Presentation")
        title_lower = title.lower()

        # Theme color mapping (Default to NotebookLM Light Blueprint #FAF9F6 unless Dark Theme requested)
        theme_style = str(spec.get("theme_style", "notebooklm")).lower()

        if theme_style in ["dark", "dark_blueprint", "void"]:
            bg_color = RGBColor(15, 23, 42)         # Dark Navy Blueprint #0F172A
            header_color = RGBColor(248, 250, 252)  # Crisp White #F8FAFC
            subtitle_color = RGBColor(148, 163, 184)# Slate Gray
            card_color = RGBColor(30, 41, 59)      # Card Fill #1E293B
            card_border = RGBColor(79, 70, 229)    # Indigo Accent #4F46E5
            text_main = RGBColor(241, 245, 249)
            accent_blue = RGBColor(96, 165, 250)
        else:
            # NotebookLM Technical Blueprint Style (Light Paper #FAF9F6)
            bg_color = RGBColor(250, 249, 246)      # Off-white Blueprint #FAF9F6
            header_color = RGBColor(17, 24, 39)     # Dark Charcoal #111827
            subtitle_color = RGBColor(75, 85, 99)   # Slate #4B5563
            card_color = RGBColor(255, 255, 255)    # Pure White Card #FFFFFF
            card_border = RGBColor(209, 213, 219)  # Tech Border #D1D5DB
            text_main = RGBColor(17, 24, 39)
            accent_blue = RGBColor(79, 70, 229)     # Indigo Accent #4F46E5

        slides_data = spec.get("slides", [])
        blank_slide_layout = prs.slide_layouts[6]

        for i, s in enumerate(slides_data):
            slide_type = s.get("type", "content").lower()
            image_path = s.get("image_path") or s.get("image_url")
            slide = prs.slides.add_slide(blank_slide_layout)

            # Slide Blueprint Canvas Background
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = bg_color
            bg.line.fill.background()

            # Technical Grid Crosshairs / Outer Frame
            frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(12.533), Inches(6.7))
            frame.fill.background()
            frame.line.color.rgb = card_border
            frame.line.width = Pt(1)

            if slide_type == "title" or i == 0:
                # Title Slide: NotebookLM Split Blueprint Layout
                if image_path and os.path.exists(image_path):
                    # Left side image diagram, Right side title text
                    try:
                        slide.shapes.add_picture(image_path, Inches(0.6), Inches(0.6), Inches(5.8), Inches(6.3))
                    except Exception:
                        pass

                    # Right Text Box
                    txBox = slide.shapes.add_textbox(Inches(6.7), Inches(1.2), Inches(5.8), Inches(5.0))
                    tf = txBox.text_frame
                    tf.word_wrap = True

                    p0 = tf.paragraphs[0]
                    p0.text = "SYSTEM ROOT // AGENT OS"
                    p0.font.size = Pt(11)
                    p0.font.bold = True
                    p0.font.color.rgb = accent_blue

                    p = tf.add_paragraph()
                    p.text = s.get("title", spec.get("title", "SeniorAgent Orchestrator"))
                    p.font.size = Pt(40)
                    p.font.bold = True
                    p.font.color.rgb = header_color
                    p.space_before = Pt(10)

                    subtitle_text = s.get("subtitle", spec.get("subtitle", ""))
                    if subtitle_text:
                        p2 = tf.add_paragraph()
                        p2.text = subtitle_text
                        p2.font.size = Pt(20)
                        p2.font.color.rgb = subtitle_color
                        p2.space_before = Pt(14)
                else:
                    # Full Title Layout
                    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(0.15), Inches(3.5))
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = accent_blue
                    bar.line.fill.background()

                    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(3.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True

                    p = tf.paragraphs[0]
                    p.text = s.get("title", spec.get("title", "Executive Presentation"))
                    p.font.size = Pt(44)
                    p.font.bold = True
                    p.font.color.rgb = header_color

                    subtitle_text = s.get("subtitle", spec.get("subtitle", ""))
                    if subtitle_text:
                        p2 = tf.add_paragraph()
                        p2.text = subtitle_text
                        p2.font.size = Pt(22)
                        p2.font.color.rgb = subtitle_color
                        p2.space_before = Pt(14)

            elif slide_type in ["split_image", "image"] or (image_path and os.path.exists(str(image_path))):
                # Split Image Slide Layout (Text Left / Image Right)
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(5.8), Inches(6.0))
                tf = txBox.text_frame
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.text = s.get("title", f"Slide {i+1}")
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = header_color

                subtitle_text = s.get("subtitle", "")
                if subtitle_text:
                    p2 = tf.add_paragraph()
                    p2.text = subtitle_text
                    p2.font.size = Pt(18)
                    p2.font.color.rgb = subtitle_color
                    p2.space_before = Pt(10)

                bullets = s.get("bullets", [])
                for b_text in bullets[:4]:
                    bp = tf.add_paragraph()
                    bp.text = f"•  {b_text}"
                    bp.font.size = Pt(15)
                    bp.font.color.rgb = text_main
                    bp.space_before = Pt(12)

                # Embed Image on Right Half
                if image_path and os.path.exists(image_path):
                    try:
                        # Image container border frame
                        img_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(0.8), Inches(5.8), Inches(5.8))
                        img_border.fill.background()
                        img_border.line.color.rgb = card_border
                        img_border.line.width = Pt(1.5)
                        
                        slide.shapes.add_picture(image_path, Inches(6.9), Inches(0.9), Inches(5.6), Inches(5.6))
                    except Exception as ie:
                        logger.error("Error adding picture to PPTX slide: %s", ie)

            elif slide_type in ["process", "pipeline"]:
                # Process Pipeline Layout (Horizontal Node Timeline)
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.0))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = s.get("title", "The Engine Pipeline")
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = header_color

                # Horizontal Timeline Arrow Connector
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.2), Inches(2.2), Inches(10.8), Inches(0.3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = accent_blue
                arrow.line.fill.background()

                steps = s.get("steps") or s.get("bullets") or ["The Brain", "The Backend", "The Frontend"]
                num_steps = min(len(steps), 3)
                col_width = Inches(11.733 / num_steps)

                for step_idx, step_text in enumerate(steps[:3]):
                    left_pos = Inches(0.8 + step_idx * (11.733 / num_steps))
                    
                    # Node Box
                    node_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.8), col_width - Inches(0.3), Inches(1.5))
                    node_box.fill.solid()
                    node_box.fill.fore_color.rgb = card_color
                    node_box.line.color.rgb = card_border
                    node_box.line.width = Pt(1.5)

                    ntf = node_box.text_frame
                    ntf.word_wrap = True
                    np1 = ntf.paragraphs[0]
                    np1.text = f"Step {step_idx+1}: {step_text.split(':')[0] if ':' in step_text else step_text}"
                    np1.font.size = Pt(18)
                    np1.font.bold = True
                    np1.font.color.rgb = header_color

                    # Code Callout Box Underneath
                    code_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(4.5), col_width - Inches(0.3), Inches(2.0))
                    code_box.fill.solid()
                    code_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
                    code_box.line.color.rgb = card_border
                    code_box.line.width = Pt(1)

                    ctf = code_box.text_frame
                    ctf.word_wrap = True
                    cp1 = ctf.paragraphs[0]
                    cp1.text = step_text.split(':')[1] if ':' in step_text else f"// Execution Config {step_idx+1}"
                    cp1.font.size = Pt(13)
                    cp1.font.color.rgb = RGBColor(51, 65, 85)

            else:
                # Standard Content Slide
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.0))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = s.get("title", f"Slide {i+1}")
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = header_color

                bullets = s.get("bullets", [])
                metrics = s.get("metrics", [])

                if metrics:
                    num_metrics = min(len(metrics), 4)
                    card_w = Inches(11.733 / max(num_metrics, 1) - 0.25)
                    for m_idx, m in enumerate(metrics[:4]):
                        left_pos = Inches(0.8 + m_idx * (11.733 / num_metrics))
                        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), card_w, Inches(2.0))
                        card.fill.solid()
                        card.fill.fore_color.rgb = card_color
                        card.line.color.rgb = card_border
                        card.line.width = Pt(1.5)

                        ctf = card.text_frame
                        ctf.word_wrap = True
                        cp1 = ctf.paragraphs[0]
                        cp1.text = str(m.get("label", "Metric")).upper()
                        cp1.font.size = Pt(11)
                        cp1.font.bold = True
                        cp1.font.color.rgb = subtitle_color

                        cp2 = ctf.add_paragraph()
                        cp2.text = str(m.get("value", "0"))
                        cp2.font.size = Pt(30)
                        cp2.font.bold = True
                        cp2.font.color.rgb = accent_blue

                if bullets:
                    top_pos = Inches(4.0) if metrics else Inches(1.8)
                    for b_idx, b_text in enumerate(bullets[:5]):
                        c_top = top_pos + Inches(b_idx * 1.0)
                        card_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), c_top, Inches(11.733), Inches(0.85))
                        card_b.fill.solid()
                        card_b.fill.fore_color.rgb = card_color
                        card_b.line.color.rgb = card_border
                        card_b.line.width = Pt(1)

                        btf = card_b.text_frame
                        btf.word_wrap = True
                        bp = btf.paragraphs[0]
                        bp.text = f"•   {b_text}"
                        bp.font.size = Pt(16)
                        bp.font.color.rgb = text_main

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)
        return {
            "status": "success",
            "file_path": output_path,
            "slides_count": len(slides_data),
            "message": f"Successfully created PowerPoint presentation at '{output_path}'"
        }
    except Exception as e:
        logger.error("Error creating PowerPoint deck: %s", e)
        return {"status": "error", "error": str(e)}


def create_html_presentation(spec: Dict[str, Any], output_path: str) -> Dict[str, Any]:
    """
    Build a NotebookLM Technical Blueprint HTML Reveal.js presentation with responsive split layouts,
    process pipelines, embedded AI images, and clean technical styling.
    """
    try:
        title = spec.get("title", "Executive Presentation")
        subtitle = spec.get("subtitle", "Technical Blueprint Deck")
        author = spec.get("author", "Business Agent")
        slides_data = spec.get("slides", [])

        html_slides = []
        for i, s in enumerate(slides_data):
            stype = s.get("type", "content").lower()
            stitle = s.get("title", f"Slide {i+1}")
            ssub = s.get("subtitle", "")
            bullets = s.get("bullets", [])
            metrics = s.get("metrics", [])
            image_path = s.get("image_path") or s.get("image_url")

            # Convert local absolute Windows image path to file URL
            img_src = ""
            if image_path:
                norm_p = str(image_path).replace("\\", "/")
                img_src = f"file:///{norm_p}" if not norm_p.startswith("http") else norm_p

            if stype == "title" or i == 0:
                if img_src:
                    slide_html = f"""
                    <section class="title-slide split-title">
                        <div class="left-col">
                            <div class="blueprint-frame">
                                <img src="{img_src}" alt="System Diagram" class="title-diagram-img" />
                            </div>
                        </div>
                        <div class="right-col">
                            <div class="tech-badge">SYSTEM ROOT // AGENT OS</div>
                            <h1 class="slide-main-title">{stitle}</h1>
                            <h3 class="slide-subtitle">{ssub or subtitle}</h3>
                            <div class="title-footer">Designed for Developers, Product Strategists & Advocates</div>
                        </div>
                    </section>
                    """
                else:
                    slide_html = f"""
                    <section class="title-slide">
                        <div class="tech-badge">SYSTEM ROOT // AGENT OS</div>
                        <h1 class="slide-main-title">{stitle}</h1>
                        <h3 class="slide-subtitle">{ssub or subtitle}</h3>
                        <div class="title-footer">Prepared by <strong>{author}</strong> | NotebookLM Blueprint Suite</div>
                    </section>
                    """
            elif stype in ["split_image", "image"] or img_src:
                b_items = "".join([f'<div class="bullet-card"><span>✦</span> {b}</div>' for b in bullets[:4]])
                slide_html = f"""
                <section class="split-slide-container">
                    <div class="split-left">
                        <h2 class="slide-header">{stitle}</h2>
                        {f'<p class="slide-subhead">{ssub}</p>' if ssub else ''}
                        <div class="bullets-wrapper">{b_items}</div>
                    </div>
                    <div class="split-right">
                        <div class="blueprint-image-box">
                            <img src="{img_src}" alt="{stitle}" class="blueprint-img" />
                            <div class="image-annotation-overlay">HIGH-RES AI VISUAL RENDER</div>
                        </div>
                    </div>
                </section>
                """
            elif stype in ["process", "pipeline"]:
                steps = s.get("steps") or bullets or ["The Brain", "The Backend", "The Frontend"]
                nodes_html = []
                for idx, st in enumerate(steps[:3]):
                    s_title = st.split(':')[0] if ':' in st else st
                    s_code = st.split(':')[1] if ':' in st else f"// Config {idx+1}"
                    nodes_html.append(f"""
                    <div class="pipeline-node">
                        <div class="node-title">Step {idx+1}: {s_title}</div>
                        <div class="code-callout-box">
                            <code>{s_code}</code>
                        </div>
                    </div>
                    """)
                slide_html = f"""
                <section class="pipeline-slide">
                    <h2 class="slide-header">{stitle}</h2>
                    <div class="pipeline-connector-line"></div>
                    <div class="pipeline-nodes-wrapper">{"".join(nodes_html)}</div>
                </section>
                """
            else:
                cards_html = ""
                if metrics:
                    cards = []
                    for m in metrics[:4]:
                        cards.append(f"""
                        <div class="metric-card">
                            <div class="metric-label">{m.get("label", "Metric")}</div>
                            <div class="metric-value">{m.get("value", "0")}</div>
                        </div>
                        """)
                    cards_html = f'<div class="metrics-grid">{"".join(cards)}</div>'

                bullets_html = ""
                if bullets:
                    b_items = "".join([f'<div class="bullet-card"><span>✦</span> {b}</div>' for b in bullets[:5]])
                    bullets_html = f'<div class="bullets-wrapper">{b_items}</div>'

                slide_html = f"""
                <section>
                    <h2 class="slide-header">{stitle}</h2>
                    {cards_html}
                    {bullets_html}
                </section>
                """

            html_slides.append(slide_html)

        reveal_template = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/theme/dracula.min.css">
    <link href="https://fonts.googleapis.com/css2?family=Outfit:wght@400;600;700;800&family=Inter:wght@400;500;600;700&display=swap" rel="stylesheet">
    <style>
        :root {{
            --bg-blueprint: #FAF9F6;
            --text-dark: #111827;
            --text-sub: #4B5563;
            --accent-indigo: #4F46E5;
            --border-tech: #D1D5DB;
            --card-bg: #FFFFFF;
        }}
        .reveal {{ font-family: 'Inter', sans-serif; background: var(--bg-blueprint); color: var(--text-dark); }}
        .reveal::before {{
            content: "";
            position: fixed;
            inset: 0;
            background-size: 30px 30px;
            background-image: linear-gradient(to right, rgba(209, 213, 219, 0.4) 1px, transparent 1px),
                              linear-gradient(to bottom, rgba(209, 213, 219, 0.4) 1px, transparent 1px);
            pointer-events: none;
        }}
        .title-slide {{ text-align: left; padding: 40px; }}
        .split-title {{ display: flex; gap: 40px; align-items: center; }}
        .left-col {{ flex: 1; }}
        .right-col {{ flex: 1.2; }}
        .tech-badge {{ display: inline-block; padding: 6px 14px; background: rgba(79, 70, 229, 0.1); border: 1px solid var(--accent-indigo); border-radius: 6px; font-size: 0.75em; font-weight: 700; color: var(--accent-indigo); letter-spacing: 1.5px; margin-bottom: 20px; }}
        .slide-main-title {{ font-family: 'Outfit', sans-serif; font-size: 3.2em !important; font-weight: 800; color: var(--text-dark); margin-bottom: 10px !important; line-height: 1.1 !important; }}
        .slide-subtitle {{ font-size: 1.4em !important; color: var(--text-sub); font-weight: 400; margin-top: 10px !important; }}
        .title-footer {{ font-size: 0.85em; color: #6B7280; margin-top: 40px; border-top: 1px solid var(--border-tech); padding-top: 15px; }}
        .slide-header {{ font-family: 'Outfit', sans-serif; text-align: left; font-size: 2.2em !important; font-weight: 800; color: var(--text-dark); margin-bottom: 25px !important; border-bottom: 2px solid var(--border-tech); padding-bottom: 12px; }}
        .split-slide-container {{ display: flex; gap: 35px; align-items: center; text-align: left; }}
        .split-left {{ flex: 1.1; }}
        .split-right {{ flex: 1.2; }}
        .blueprint-image-box {{ position: relative; border: 2px solid var(--border-tech); border-radius: 12px; padding: 8px; background: #FFFFFF; box-shadow: 0 10px 30px rgba(0,0,0,0.08); }}
        .blueprint-img {{ width: 100%; height: 420px; object-fit: cover; border-radius: 8px; }}
        .image-annotation-overlay {{ position: absolute; bottom: 16px; left: 16px; background: rgba(17, 24, 39, 0.85); color: #FFFFFF; padding: 4px 10px; font-size: 0.65em; font-weight: 700; border-radius: 4px; letter-spacing: 1px; }}
        .pipeline-slide {{ text-align: left; }}
        .pipeline-connector-line {{ height: 6px; background: var(--accent-indigo); margin: 30px 0; border-radius: 3px; position: relative; }}
        .pipeline-nodes-wrapper {{ display: grid; grid-template-columns: repeat(3, 1fr); gap: 20px; }}
        .pipeline-node {{ background: #FFFFFF; border: 1px solid var(--border-tech); border-radius: 10px; padding: 20px; box-shadow: 0 4px 12px rgba(0,0,0,0.05); }}
        .node-title {{ font-family: 'Outfit', sans-serif; font-weight: 700; font-size: 1.1em; color: var(--text-dark); margin-bottom: 12px; }}
        .code-callout-box {{ background: #F3F4F6; border: 1px solid #E5E7EB; border-radius: 6px; padding: 12px; font-family: monospace; font-size: 0.85em; color: #1F2937; }}
        .metrics-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 20px; margin: 25px 0; }}
        .metric-card {{ background: #FFFFFF; border: 1px solid var(--border-tech); border-radius: 12px; padding: 24px; text-align: left; box-shadow: 0 4px 12px rgba(0,0,0,0.05); }}
        .metric-label {{ font-size: 0.75em; text-transform: uppercase; color: var(--text-sub); letter-spacing: 1px; font-weight: 700; }}
        .metric-value {{ font-family: 'Outfit', sans-serif; font-size: 2.5em; font-weight: 800; color: var(--accent-indigo); margin: 8px 0 0 0; }}
        .bullets-wrapper {{ display: flex; flex-direction: column; gap: 12px; text-align: left; margin-top: 15px; }}
        .bullet-card {{ background: #FFFFFF; border: 1px solid var(--border-tech); border-left: 4px solid var(--accent-indigo); border-radius: 8px; padding: 16px 20px; font-size: 1.1em; color: var(--text-dark); box-shadow: 0 2px 8px rgba(0,0,0,0.04); }}
        .bullet-card span {{ color: var(--accent-indigo); margin-right: 8px; }}
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
            {"".join(html_slides)}
        </div>
    </div>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.js"></script>
    <script>
        Reveal.initialize({{
            controls: true,
            progress: true,
            center: false,
            hash: true,
            transition: 'slide'
        }});
    </script>
</body>
</html>"""

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(reveal_template)

        return {
            "status": "success",
            "file_path": output_path,
            "slides_count": len(slides_data),
            "message": f"Successfully created interactive HTML presentation at '{output_path}'"
        }
    except Exception as e:
        logger.error("Error creating HTML presentation: %s", e)
        return {"status": "error", "error": str(e)}
