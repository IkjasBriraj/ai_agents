"""
Presentation Service
Generates professional PowerPoint presentations (.pptx) and interactive HTML slide decks (Reveal.js)
from structured JSON slide specifications.
"""

import os
import json
import logging
from typing import Dict, Any, List, Optional

logger = logging.getLogger(__name__)


def create_powerpoint_deck(spec: Dict[str, Any], output_path: str) -> Dict[str, Any]:
    """
    Build a PowerPoint presentation (.pptx) from a JSON slide specification.
    
    Spec format:
    {
        "title": "Presentation Title",
        "subtitle": "Subtitle or Company Name",
        "author": "Business Agent",
        "theme_color": "#1E3A8A",  # Primary color hex
        "slides": [
            {
                "type": "title",  # title, content, metrics, table, split, closing
                "title": "Slide Title",
                "subtitle": "Optional Subtitle",
                "bullets": ["Point 1", "Point 2"],
                "metrics": [{"label": "ARR", "value": "$2.4M", "change": "+45%"}, ...],
                "table": {"headers": ["Col1", "Col2"], "rows": [["Val1", "Val2"]]}
            }
        ]
    }
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE

        prs = Presentation()
        # Set 16:9 widescreen layout
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Parse theme colors
        hex_color = spec.get("theme_color", "#0F172A").lstrip('#')
        if len(hex_color) == 6:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            brand_color = RGBColor(r, g, b)
        else:
            brand_color = RGBColor(15, 23, 42)  # Slate dark

        text_dark = RGBColor(30, 41, 59)
        text_light = RGBColor(248, 250, 252)
        accent_blue = RGBColor(37, 99, 235)

        slides_data = spec.get("slides", [])

        blank_slide_layout = prs.slide_layouts[6]  # Blank layout

        for i, s in enumerate(slides_data):
            slide_type = s.get("type", "content").lower()
            slide = prs.slides.add_slide(blank_slide_layout)

            if slide_type == "title" or i == 0:
                # Dark Brand Cover Background
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
                bg.fill.solid()
                bg.fill.fore_color.rgb = brand_color
                bg.line.fill.background()

                # Title text box
                txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
                tf = txBox.text_frame
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.text = s.get("title", spec.get("title", "Executive Presentation"))
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = text_light
                p.alignment = PP_ALIGN.LEFT

                subtitle_text = s.get("subtitle", spec.get("subtitle", ""))
                if subtitle_text:
                    p2 = tf.add_paragraph()
                    p2.text = subtitle_text
                    p2.font.size = Pt(24)
                    p2.font.color.rgb = RGBColor(148, 163, 184)
                    p2.alignment = PP_ALIGN.LEFT

                # Author / Date footer
                p3 = tf.add_paragraph()
                p3.text = f"Prepared by: {spec.get('author', 'Business Agent')} | AI Generated Strategy Deck"
                p3.font.size = Pt(14)
                p3.font.color.rgb = RGBColor(100, 116, 139)

            else:
                # Top header bar
                header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
                header_bg.fill.solid()
                header_bg.fill.fore_color.rgb = brand_color
                header_bg.line.fill.background()

                # Slide Title
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = s.get("title", f"Slide {i+1}")
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = text_light

                # Content layout
                bullets = s.get("bullets", [])
                metrics = s.get("metrics", [])
                table_data = s.get("table", {})

                if metrics:
                    # Key Metrics Grid Layout
                    num_metrics = min(len(metrics), 4)
                    card_width = Inches(11.5 / max(num_metrics, 1) - 0.3)
                    
                    for m_idx, metric in enumerate(metrics[:4]):
                        left = Inches(0.8 + m_idx * (11.5 / num_metrics))
                        top = Inches(1.6)

                        # Card Box
                        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, Inches(2.2))
                        card.fill.solid()
                        card.fill.fore_color.rgb = RGBColor(241, 245, 249)
                        card.line.color.rgb = RGBColor(203, 213, 225)

                        ctf = card.text_frame
                        ctf.word_wrap = True
                        
                        cp1 = ctf.paragraphs[0]
                        cp1.text = str(metric.get("label", "Metric")).upper()
                        cp1.font.size = Pt(12)
                        cp1.font.bold = True
                        cp1.font.color.rgb = RGBColor(100, 116, 139)

                        cp2 = ctf.add_paragraph()
                        cp2.text = str(metric.get("value", "0"))
                        cp2.font.size = Pt(36)
                        cp2.font.bold = True
                        cp2.font.color.rgb = accent_blue

                        if metric.get("change"):
                            cp3 = ctf.add_paragraph()
                            cp3.text = str(metric.get("change"))
                            cp3.font.size = Pt(14)
                            cp3.font.bold = True
                            cp3.font.color.rgb = RGBColor(16, 185, 129)  # Emerald green

                if bullets:
                    top_pos = Inches(4.2) if metrics else Inches(1.6)
                    bullet_box = slide.shapes.add_textbox(Inches(0.8), top_pos, Inches(11.533), Inches(5.0))
                    btf = bullet_box.text_frame
                    btf.word_wrap = True

                    for b_idx, b_text in enumerate(bullets):
                        bp = btf.paragraphs[0] if b_idx == 0 else btf.add_paragraph()
                        bp.text = f"•  {b_text}"
                        bp.font.size = Pt(18)
                        bp.font.color.rgb = text_dark
                        bp.space_after = Pt(12)

                if table_data and "headers" in table_data and "rows" in table_data:
                    headers = table_data["headers"]
                    rows = table_data["rows"]
                    
                    rows_count = len(rows) + 1
                    cols_count = len(headers)
                    
                    left = Inches(0.8)
                    top = Inches(2.0)
                    width = Inches(11.733)
                    height = Inches(0.5 * rows_count)

                    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
                    tbl = table_shape.table

                    # Header Row
                    for c_idx, head in enumerate(headers):
                        cell = tbl.cell(0, c_idx)
                        cell.text = str(head)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = brand_color
                        for p in cell.text_frame.paragraphs:
                            p.font.bold = True
                            p.font.color.rgb = text_light
                            p.font.size = Pt(14)

                    # Data Rows
                    for r_idx, row in enumerate(rows):
                        for c_idx, val in enumerate(row):
                            cell = tbl.cell(r_idx + 1, c_idx)
                            cell.text = str(val)
                            cell.fill.solid()
                            if r_idx % 2 == 0:
                                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)
                            else:
                                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            for p in cell.text_frame.paragraphs:
                                p.font.size = Pt(13)
                                p.font.color.rgb = text_dark

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)
        
        return {
            "status": "success",
            "file_path": output_path,
            "slides_count": len(slides_data),
            "message": f"Successfully created PowerPoint presentation at '{output_path}'"
        }
    except ImportError:
        logger.warning("python-pptx is not installed.")
        return {
            "status": "error",
            "error": "python-pptx library is not installed. Install via pip install python-pptx."
        }
    except Exception as e:
        logger.error("Error creating PowerPoint deck: %s", e)
        return {
            "status": "error",
            "error": str(e)
        }


def create_html_presentation(spec: Dict[str, Any], output_path: str) -> Dict[str, Any]:
    """
    Build an interactive, single-file HTML Reveal.js presentation.
    Can be viewed directly in web browser or previewed in frontend modals.
    """
    try:
        title = spec.get("title", "Executive Presentation")
        subtitle = spec.get("subtitle", "Business Strategy Deck")
        author = spec.get("author", "Business Agent")
        slides_data = spec.get("slides", [])

        html_slides = []

        for i, s in enumerate(slides_data):
            stype = s.get("type", "content").lower()
            stitle = s.get("title", f"Slide {i+1}")
            ssub = s.get("subtitle", "")
            bullets = s.get("bullets", [])
            metrics = s.get("metrics", [])
            table_data = s.get("table", {})

            if stype == "title" or i == 0:
                slide_html = f"""
                <section class="title-slide">
                    <h1 class="slide-main-title">{stitle}</h1>
                    <h3 class="slide-subtitle">{ssub or subtitle}</h3>
                    <div class="title-footer">Prepared by: <strong>{author}</strong> | AI Generated Strategy Deck</div>
                </section>
                """
            else:
                cards_html = ""
                if metrics:
                    cards = []
                    for m in metrics:
                        chg_html = f'<div class="metric-change">{m.get("change")}</div>' if m.get("change") else ''
                        cards.append(f"""
                        <div class="metric-card">
                            <div class="metric-label">{m.get("label", "Metric")}</div>
                            <div class="metric-value">{m.get("value", "0")}</div>
                            {chg_html}
                        </div>
                        """)
                    cards_html = f'<div class="metrics-grid">{"".join(cards)}</div>'

                bullets_html = ""
                if bullets:
                    b_items = "".join([f'<li>{b}</li>' for b in bullets])
                    bullets_html = f'<ul class="bullets-list">{b_items}</ul>'

                table_html = ""
                if table_data and "headers" in table_data and "rows" in table_data:
                    ths = "".join([f'<th>{h}</th>' for h in table_data["headers"]])
                    trs = []
                    for row in table_data["rows"]:
                        tds = "".join([f'<td>{v}</td>' for v in row])
                        trs.append(f'<tr>{tds}</tr>')
                    table_html = f'<table class="slide-table"><thead><tr>{ths}</tr></thead><tbody>{"".join(trs)}</tbody></table>'

                slide_html = f"""
                <section>
                    <h2 class="slide-header">{stitle}</h2>
                    {cards_html}
                    {bullets_html}
                    {table_html}
                </section>
                """

            html_slides.append(slide_html)

        reveal_template = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/theme/dracula.min.css">
    <style>
        .reveal {{ font-family: 'Inter', system-ui, sans-serif; }}
        .title-slide {{ text-align: left; padding: 40px; }}
        .slide-main-title {{ font-size: 2.8em !important; font-weight: 800; background: linear-gradient(135deg, #60A5FA, #3B82F6); -webkit-background-clip: text; -webkit-text-fill-color: transparent; }}
        .slide-subtitle {{ font-size: 1.4em !important; color: #94A3B8; font-weight: 300; margin-top: 15px !important; }}
        .title-footer {{ font-size: 0.8em; color: #64748B; margin-top: 50px; border-t: 1px solid #334155; padding-top: 15px; }}
        .slide-header {{ text-align: left; font-size: 1.8em !important; font-weight: 700; color: #F8FAFC; margin-bottom: 30px !important; border-bottom: 2px solid #3B82F6; padding-bottom: 10px; }}
        .metrics-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 20px; margin: 30px 0; }}
        .metric-card {{ background: #1E293B; border: 1px solid #334155; border-radius: 8px; padding: 20px; text-align: left; }}
        .metric-label {{ font-size: 0.7em; text-transform: uppercase; color: #94A3B8; letter-spacing: 1px; font-weight: 700; }}
        .metric-value {{ font-size: 2.2em; font-weight: 800; color: #60A5FA; margin: 10px 0 5px 0; }}
        .metric-change {{ font-size: 0.8em; font-weight: 700; color: #10B981; }}
        .bullets-list {{ text-align: left; font-size: 1.1em; line-height: 1.8; color: #E2E8F0; margin-left: 30px; }}
        .bullets-list li {{ margin-bottom: 12px; }}
        .slide-table {{ width: 100%; border-collapse: collapse; font-size: 0.85em; margin-top: 20px; }}
        .slide-table th {{ background: #1E3A8A; color: #FFFFFF; padding: 12px; text-align: left; border: 1px solid #3B82F6; }}
        .slide-table td {{ padding: 10px; border: 1px solid #334155; background: #0F172A; color: #E2E8F0; }}
        .slide-table tr:nth-child(even) td {{ background: #1E293B; }}
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
            {"".join(html_slides)}
        </div>
    </div>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.js"></script>
    <script>
        Reveal.initialize({{
            controls: true,
            progress: true,
            center: false,
            hash: true,
            transition: 'slide'
        }});
    </script>
</body>
</html>"""

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(reveal_template)

        return {
            "status": "success",
            "file_path": output_path,
            "slides_count": len(slides_data),
            "message": f"Successfully created interactive HTML presentation at '{output_path}'"
        }
    except Exception as e:
        logger.error("Error creating HTML presentation: %s", e)
        return {
            "status": "error",
            "error": str(e)
        }
